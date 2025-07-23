import os
import subprocess
from threading import Thread

from evdev import InputDevice, list_devices, categorize

from kivy.app import App
from kivy.core.window import Window
from kivy.properties import ObjectProperty, StringProperty, BooleanProperty
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.floatlayout import FloatLayout
from kivy.uix.gridlayout import GridLayout
from kivy.uix.label import Label
from kivy.clock import Clock

from pdf2image import convert_from_bytes
from pptx import Presentation

FOOT_SWITCH_DEVICE_NAME_SUFFIX = "FootSwitch Keyboard"
FOOT_SWITCH_DEVICE_A_KEY = "KEY_A"
FOOT_SWITCH_DEVICE_B_KEY = "KEY_B"
FOOT_SWITCH_DEVICE_C_KEY = "KEY_C"
SONGBOOK_MIN_ROWS_NUM = 3
SONGBOOK_MIN_COLS_NUM = 6
SONGBOOKS_FOLDER = "songbooks"  # by default one folder up from code folder  otherwise within code folder
TEMP_FOLDER = "#converted#"  # inside SONGBOOKS_FOLDER

# SIZES
TOP_BAR_TO_IMAGE_RATIO = 0.04
BOTTOM_BAR_TO_IMAGE_RATIO = 0.04


class LoadingScreenLayout(BoxLayout):
    previous_text = ""

    def draw_text(self, text, append=False):
        
        if append:
            all = self.previous_text.split("\n")
            if len(all) > 0:
                all[-1] = all[-1] + text
            else:
                all.append(text)
            self.loading_screen_text = "\n".join(all)
            self.previous_text = "\n".join(all)
        else:
            self.loading_screen_text = self.previous_text + "\n" + text
            self.previous_text = self.previous_text + "\n" + text


class HomeLayout(BoxLayout):
    pass


class Songbook(BoxLayout):
    focus = ObjectProperty()
    index = ObjectProperty()
    sequence = StringProperty()
    title = StringProperty()
    songs = ObjectProperty()


class SongbookLayoutMain(BoxLayout):
    pass


class BackButton(
    BoxLayout,
):
    focus = ObjectProperty()


class SongList(GridLayout):
    pass


class Song(
    BoxLayout,
):
    focus = ObjectProperty()
    is_placeholder = BooleanProperty()
    sequence = ObjectProperty()
    artist = ObjectProperty()
    song = ObjectProperty()

    def __init__(
        self,
        sequence=None,
        artist=None,
        song=None,
        index=None,
        is_placeholder=False,
        images=None,
        **kwargs,
    ):
        super().__init__(**kwargs)
        self.focus = False
        self.index = index
        self.is_placeholder = is_placeholder
        self.images = images
        self.sequence = str(index + 1)  # I decided to do my own counting and not show the users sequence identifier
        self.artist = artist
        self.song = song


class SequenceLabel(Label):
    pass


class ArtistLabel(Label):
    pass


class SongLabel(Label):
    pass


class PromptLayout(BoxLayout):
    current_image_number = ObjectProperty(0)
    current_image_source = ObjectProperty()
    number_of_slides = ObjectProperty()
    
    current_song = ObjectProperty()
    next_song = ObjectProperty()
    
    all_songs = ObjectProperty()
    placeholders_num = ObjectProperty()

    def load(self, current_song):
        
        # Prepare first draw
        self.current_song = current_song
        self.images = current_song.images
        self.number_of_slides = len(current_song.images)
        self.current_image_number = 0
        self.current_image_source = self.images[0]
        self.next_song = self.get_next_song()
    
    def get_previous_song(self):
        next_index = self.current_song.index - 1
        if next_index < 0:
            next_index = len(self.all_songs) - 1 - self.placeholders_num
        for song in self.all_songs:
            if song.index == next_index:
                return song
        return None

    def get_next_song(self):
        next_index = self.current_song.index + 1
        if next_index >= len(self.all_songs) - self.placeholders_num:
            next_index = 0
        for song in self.all_songs:
            if song.index == next_index:
                return song
        return None
    
    def prev_image(self):
        if self.current_image_number - 1 < 0:
            to_load = self.get_previous_song()
            self.load(to_load)
            self.current_image_number = self.number_of_slides - 1
            self.current_image_source = self.images[self.current_image_number]
        else:
            self.current_image_number = self.current_image_number - 1
            self.current_image_source = self.images[self.current_image_number]

    def next_image(self):
        if self.current_image_number + 1 > self.number_of_slides - 1:
            to_load = self.get_next_song()
            self.load(to_load)
        else:
            self.current_image_number = self.current_image_number + 1
            self.current_image_source = self.images[self.current_image_number]


class PrompterTopBar(BoxLayout):
    pass


class PrompterBottomBar(BoxLayout):
    pass


class TeleprompterWidget(FloatLayout):

    mode = StringProperty("loading")  # Expose to template
    current_songbook = ObjectProperty()
    focused_songbook = ObjectProperty()
    focused_song = ObjectProperty()

    def __init__(self, **kwargs):
        super().__init__(**kwargs)

        self._input_state = None
        self._song_instances = None
        self._placeholders_num = 0

        # Check for Foot Switch and connect
        self._fs_device = self._find_foot_switch_device()
        Thread(target=self._detect_foot_switch_events, daemon=True).start()

        # Connect Keyboard
        self._keyboard = Window.request_keyboard(self._keyboard_closed, self)
        self._keyboard.bind(
            on_key_down=self._on_keyboard_down,
            on_key_up=self._on_keyboard_up,
        )
        self._previous_keyboard_state = None
        
        # Define location of songbooks folder
        potential_folder = os.path.join(os.path.dirname(os.path.dirname(os.path.realpath(__file__))), SONGBOOKS_FOLDER)
        if os.path.exists(potential_folder):
            self.songbooks_path = potential_folder
            self.songbooks_converted_path = os.path.join(potential_folder, TEMP_FOLDER)
        else:
            potential_folder = os.path.join(os.path.dirname(os.path.realpath(__file__)), SONGBOOKS_FOLDER)
            if os.path.exists(potential_folder):
                self.songbooks_path = potential_folder
                self.songbooks_converted_path = os.path.join(potential_folder, TEMP_FOLDER)
            else:
                raise Exception("Could not find a songbooks folder inside of code folder or one level up")
        
        # Create a thread to load songbooks and draw its contents
        self.songbooks = []
        Thread(target=self.load_and_draw).start()

    """
    Setup UI
    """

    def _find_foot_switch_device(self):
        fs_device = None
        devices = [InputDevice(path) for path in list_devices()]

        for device in devices:
            if device.name.endswith(FOOT_SWITCH_DEVICE_NAME_SUFFIX):
                message = f"Found {FOOT_SWITCH_DEVICE_NAME_SUFFIX} at {device.path}"
                self.update_loading_screen(message)
                fs_device = device
                break
        if fs_device is None:
            self.update_loading_screen(f"Did not find footswitch device called '{FOOT_SWITCH_DEVICE_NAME_SUFFIX}'")
            self.update_loading_screen("Use keyboard instead")
            if len(devices) > 0:
                print("Found:")
                for device in devices:
                    print(device)
        return fs_device

    def _detect_foot_switch_events(self):
        if self._fs_device:
            with self._fs_device.grab_context():
                for ev in self._fs_device.async_read_loop():
                    btn = None
                    state = None
                    as_string = str(categorize(ev))
                    if FOOT_SWITCH_DEVICE_A_KEY in as_string:
                        btn = "A"
                    elif FOOT_SWITCH_DEVICE_B_KEY in as_string:
                        btn = "B"
                    elif FOOT_SWITCH_DEVICE_C_KEY in as_string:
                        btn = "C"

                    if "up" in as_string:
                        state = "up"
                    elif "down" in as_string:
                        state = "down"
                    elif "hold" in as_string:
                        state = "hold"

                    if btn and state:
                        self._input_state = (btn, state)
                        
                        # The foot switch detection is in a separate thread. Use Clock to get back into the kivy thread
                        Clock.schedule_once(lambda dt: self._decide_action())

    def _keyboard_closed(self):
        # Do not unbind, otherwise escape from prompt will switch to home but then no keys are detected anymore
        pass
        # self._keyboard.unbind(on_key_down=self._on_keyboard_down)
        # self._keyboard = None

    def _on_keyboard_down(self, keyboard, keycode, text, modifiers):
        key = keycode[1]
        btn = None
        state = None
        if key == "1" or key == "numpad1":
            btn = "A"
            if self._previous_keyboard_state == "down" + btn:
                state = "hold"
            else:
                state = "down"
        if key == "2" or key == "numpad2":
            btn = "B"
            if self._previous_keyboard_state == "down" + btn:
                state = "hold"
            else:
                state = "down"
        if key == "3" or key == "numpad3":
            btn = "C"
            if self._previous_keyboard_state == "down" + btn:
                state = "hold"
            else:
                state = "down"
        if key == "escape":
            # Stop listener
            if self.mode == "home":
                keyboard.release()
                App.get_running_app().stop()
            elif self.mode == "songbook":
                self.set_mode("home")
                self._keyboard.bind(
                    on_key_down=self._on_keyboard_down,
                    on_key_up=self._on_keyboard_up,
                )
            else:
                self.set_mode("songbook")
                self._keyboard.bind(
                    on_key_down=self._on_keyboard_down,
                    on_key_up=self._on_keyboard_up,
                )

        if btn and state:
            self._input_state = (btn, state)
            self._previous_keyboard_state = state + btn
            self._decide_action()
        return True

    def _on_keyboard_up(self, keyboard, keycode):
        key = keycode[1]
        btn = None
        state = None
        if key == "1" or key == "numapd1":
            btn = "A"
            state = "up"
        if key == "2" or key == "numapd2":
            btn = "V"
            state = "up"
        if key == "3" or key == "numapd3":
            btn = "C"
            state = "up"
        if btn and state:
            self._input_state = (btn, state)
            # self._decide_action()
        return True

    def _decide_action(self):
        if self._input_state:
            if self.mode == "home":
                if self._input_state[0] == "A" and self._input_state[1] in [
                    "hold",
                    "down",
                ]:
                    self.focus_previous_songbook()
                if self._input_state[0] == "C" and self._input_state[1] in [
                    "hold",
                    "down",
                ]:
                    self.focus_next_songbook()

                if self._input_state[0] == "B" and self._input_state[1] in [
                    "hold",
                    "down",
                ]:
                    self.songbook_open(self.focused_songbook)
            elif self.mode == "songbook":
                if self._input_state[0] == "A" and self._input_state[1] in [
                    "hold",
                    "down",
                ]:
                    self.focus_previous_song()
                if self._input_state[0] == "C" and self._input_state[1] in [
                    "hold",
                    "down",
                ]:
                    self.focus_next_song()

                if self._input_state[0] == "B" and self._input_state[1] in [
                    "hold",
                    "down",
                ]:
                    self.enter_prompt()
            else:
                if self._input_state[0] == "A" and self._input_state[1] in [
                    "hold",
                    "down",
                ]:
                    self.prompt_prev()
                if self._input_state[0] == "C" and self._input_state[1] in [
                    "hold",
                    "down",
                ]:
                    self.prompt_next()

                if self._input_state[0] == "B" and self._input_state[1] in [
                    "hold",
                    "down",
                ]:
                    self.set_mode("songbook")

    """
    Actions
    """

    def focus_previous_song(self):
        
        # Handle back button
        if self.ids["back_button"].focus is True:
            next_index = len(self._song_instances) - 1 - self._placeholders_num
            self.ids["back_button"].focus = False
        else:
            next_index = self.focused_song.index - 1    

        if next_index < 0:
            self.ids["back_button"].focus = True
        for song in self._song_instances:
            if song.index == next_index:
                song.focus = True
                self.focused_song = song
            else:
                song.focus = False

    def focus_next_song(self):

        # Handle back button
        if self.ids["back_button"].focus is True:
            next_index = 0
            self.ids["back_button"].focus = False
        else:
            next_index = self.focused_song.index + 1

        if next_index >= len(self._song_instances) - self._placeholders_num:
            self.ids["back_button"].focus = True
        for song in self._song_instances:
            if song.index == next_index:
                song.focus = True
                self.focused_song = song
            else:
                song.focus = False

    def focus_previous_songbook(self):
        next_index = self.focused_songbook.index - 1
        if next_index < 0:
            next_index = len(self.songbooks) - 1
        for songbook in self.songbooks:
            if songbook.index == next_index:
                songbook.focus = True
                self.focused_songbook = songbook
            else:
                songbook.focus = False

    def focus_next_songbook(self):
        next_index = self.focused_songbook.index + 1
        if next_index >= len(self.songbooks):
            next_index = 0
        for songbook in self.songbooks:
            if songbook.index == next_index:
                songbook.focus = True
                self.focused_songbook = songbook
            else:
                songbook.focus = False

    def enter_prompt(self):
        if self.ids["back_button"].focus is True:
            self.set_mode("home")
            return
        
        self.ids["prompt_layout"].load(self.focused_song)
        self.set_mode("prompt")

    def prompt_prev(self):
        self.ids["prompt_layout"].prev_image()

    def prompt_next(self):
        self.ids["prompt_layout"].next_image()

    def songbook_open(self, songbook):
        self.ids["back_button"].focus = False
        self.initialize_songbook(songbook)
        self.set_mode("songbook")

    def set_mode(self, mode):
        self.mode = mode

    """
    Load Song Books
    """

    def _presentation_to_images(self, path_to_presentation, songbook_name):
        IMAGE_FORMAT = "jpg"
        
        converted_path = os.path.join(self.songbooks_converted_path, songbook_name)
        if not os.path.exists(converted_path):
            os.makedirs(converted_path)

        filename_bare = os.path.basename(path_to_presentation).replace(".pptx", "")

        # Give cache if images are present but only if ppt is not newer
        expected_image_paths = []
        cache_ok = True
        try:
            presentation = Presentation(path_to_presentation)
            for i in range(len(presentation.slides)):
                expected_image = os.path.join(
                    converted_path, f"{filename_bare}-{i}.{IMAGE_FORMAT}"
                )
                expected_image_paths.append(expected_image)
                if os.path.exists(expected_image) is False:
                    cache_ok = False
                else:
                    if os.path.getmtime(expected_image) < os.path.getmtime(path_to_presentation):
                        cache_ok = False
        except Exception as e:
            print("Error during presentation loading: {}".format(e))
            cache_ok = False

        if cache_ok:
            self.update_loading_screen(" (cache)", append=True)
            return expected_image_paths

        # convert pptx to PDF
        command_list = ["soffice", "--headless", "--convert-to", "pdf", path_to_presentation]
        subprocess.run(command_list)

        pdffile_name = filename_bare + ".pdf"
        with open(pdffile_name, "rb") as f:
            pdf_bytes = f.read()
        images = convert_from_bytes(pdf_bytes, dpi=96 * 4)

        created_image_paths = []
        for i, img in enumerate(images):
            im_name = os.path.join(converted_path, f"{filename_bare}-{i}.{IMAGE_FORMAT}")
            created_image_paths.append(im_name)
            img.save(im_name)

        os.unlink(pdffile_name)

        self.update_loading_screen(" (converted)", append=True)

        return created_image_paths
    
    def load_songbooks(self):
        
        # Get candidates
        potential_songbooks = []
        for songbook_folder in sorted(os.listdir(self.songbooks_path)):
            songbook_path = os.path.join(self.songbooks_path, songbook_folder)
            if os.path.isdir(songbook_path) and songbook_folder != TEMP_FOLDER and "-" in songbook_folder:
                potential_songbooks.append(songbook_path)

        songbooks = []
        for songbook_index, songbook_path in enumerate(potential_songbooks):
            songbook_folder = os.path.basename(songbook_path)
            songbook_sequence = songbook_folder.split("-")[0].strip()
            songbook_title = songbook_folder.split("-")[1].strip()
            
            self.update_loading_screen("")
            self.update_loading_screen("SONGBOOK: " + songbook_title)


            # Collect songs for this songbook
            songs = []
            for f in sorted(os.listdir(songbook_path)):
                if f.startswith("~"):
                    continue
                if f.endswith("pptx"):
                    info = f.replace(".pptx", "")
                    sequence = info.split("-")[0].strip()
                    artist = info.split("-")[1].strip()
                    song = info.split("-")[2].strip()
                    
                    self.update_loading_screen(f"{artist} - {song}")
                    song = {
                        "sequence": sequence,
                        "artist": artist,
                        "song": song,
                        "images": self._presentation_to_images(
                            os.path.join(songbook_path, f),
                            songbook_folder
                        )
                    }
                    songs.append(song)

            # Collect songbooks
            songbooks.append({
                "sequence": songbook_sequence,
                "title": songbook_title,
                "songs": songs,
                "index": songbook_index
            })
        return songbooks
    
    def update_loading_screen(self, message, append=False):
        self.ids["loading_screen"].draw_text(message, append=append)

    def initialize_home(self):

        # Create song widgets
        for songbook in self.songbooks:
            self.ids["home_layout"].add_widget(songbook)

        # Focus first
        self.songbooks[0].focus = True
        self.focused_songbook = self.songbooks[0]

    def initialize_songbook(self, songbook):
        
        # Reset
        self.ids["song_list"].clear_widgets()
        self._placeholders_num = 0
        self.ids["song_list"].rows = SONGBOOK_MIN_ROWS_NUM
        self.ids["song_list"].cols = SONGBOOK_MIN_COLS_NUM
        
        # Create song widgets
        self._song_instances = []
        for index, c in enumerate(songbook.songs):
            song_instance = Song(
                sequence=c["sequence"],
                artist=c["artist"],
                song=c["song"],
                images=c.get("images", []),
                is_placeholder=c.get("is_placeholder", False),
                index=index,
            )
            self._song_instances.append(song_instance)
            self.ids["song_list"].add_widget(song_instance)

        # Do we need placeholders ?
        if len(songbook.songs) < SONGBOOK_MIN_COLS_NUM * SONGBOOK_MIN_ROWS_NUM:
            to_add = SONGBOOK_MIN_COLS_NUM * SONGBOOK_MIN_ROWS_NUM - len(songbook.songs)
            self._placeholders_num = to_add
            while to_add > 0:
                placeholder_song_instance = (
                    Song(
                        is_placeholder=True,
                        sequence="",
                        artist="",
                        song="",
                        index=len(songbook.songs) + to_add
                    )
                )
                self._song_instances.append(placeholder_song_instance)
                self.ids["song_list"].add_widget(placeholder_song_instance)
                to_add = to_add - 1
        else:
            # Need to increase grid space
            rows_needed = len(songbook.songs) // SONGBOOK_MIN_COLS_NUM
            self.ids["song_list"].rows = rows_needed + 1


        # Focus first song
        self._song_instances[0].focus = True
        self.focused_song = self._song_instances[0]
        
        self.ids["prompt_layout"].all_songs = self._song_instances
        self.ids["prompt_layout"].placeholders_num = self._placeholders_num
    
    def load_and_draw(self):

        def _draw(_songbooks):
            for songbook_dict in _songbooks:
                self.songbooks.append(Songbook(
                    sequence=songbook_dict["sequence"],
                    title=songbook_dict["title"],
                    songs=songbook_dict["songs"],
                    index=songbook_dict["index"],
                    focus=False
                ))
            self.current_songbook = self.songbooks[0] if self.songbooks else None
            self.initialize_home()
            self.set_mode("home")
        
        # Use Clock otherwise we will not be in the Kivy thread at that point since
        # the whole load_and_draw method is in a separate thread
        self.update_loading_screen("\nLoading songbooks ...")
        songbook_dicts = self.load_songbooks()
        Clock.schedule_once(lambda dt: _draw(songbook_dicts))


class TeleprompterApp(App):
    
    TOP_BAR_TO_IMAGE_RATIO = TOP_BAR_TO_IMAGE_RATIO
    BOTTOM_BAR_TO_IMAGE_RATIO = BOTTOM_BAR_TO_IMAGE_RATIO
    
    Window.fullscreen = True
    Window.allow_screensaver = False

    def build(self):
        main = TeleprompterWidget()
        return main


if __name__ == "__main__":
    TeleprompterApp().run()
